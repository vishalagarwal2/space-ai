# knowledge_base/parsers.py

import io
import json
import os
import tempfile
from abc import ABC, abstractmethod
from typing import Dict, List, Optional, Any, Tuple
import logging
from pdf2image import convert_from_path 
import base64

# Document processing
import PyPDF2
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
from bs4 import BeautifulSoup

# Image processing
from PIL import Image, ImageEnhance, ImageFilter
import pytesseract
try:
    import easyocr
    EASYOCR_AVAILABLE = True
except ImportError:
    EASYOCR_AVAILABLE = False

# Audio/Video processing
try:
    import whisper
    from pydub import AudioSegment
    import cv2
    MEDIA_PROCESSING_AVAILABLE = True
except ImportError:
    MEDIA_PROCESSING_AVAILABLE = False

# Text processing
import re
from collections import defaultdict
from mistralai import Mistral, DocumentURLChunk, ImageURLChunk, TextChunk
from pathlib import Path
from mistralai.models import OCRResponse

logger = logging.getLogger(__name__)


class BaseParser(ABC):
    """Base class for all document parsers"""
    
    def __init__(self, config: Dict[str, Any] = None):
        self.config = config or {}
    
    @abstractmethod
    def can_parse(self, file_path: str, mime_type: str) -> bool:
        """Check if this parser can handle the given file"""
        pass
    
    @abstractmethod
    def parse(self, file_path: str, **kwargs) -> Dict[str, Any]:
        """Parse the file and return extracted content"""
        pass
    
    def get_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from the file"""
        try:
            stat = os.stat(file_path)
            return {
                'file_size': stat.st_size,
                'modified_time': stat.st_mtime,
                'parser_used': self.__class__.__name__,
            }
        except Exception as e:
            logger.error(f"Error getting metadata for {file_path}: {e}")
            return {}


class TextParser(BaseParser):
    """Parser for plain text files"""
    
    SUPPORTED_EXTENSIONS = ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.py', '.js', '.css']
    SUPPORTED_MIME_TYPES = [
        'text/plain', 'text/markdown', 'text/csv', 'application/json',
        'text/xml', 'text/html', 'text/javascript', 'text/css'
    ]
    
    def can_parse(self, file_path: str, mime_type: str) -> bool:
        ext = os.path.splitext(file_path)[1].lower()
        return ext in self.SUPPORTED_EXTENSIONS or mime_type in self.SUPPORTED_MIME_TYPES
    
    def parse(self, file_path: str, **kwargs) -> Dict[str, Any]:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Special handling for different text types
            if file_path.endswith('.json'):
                content = self._format_json(content)
            elif file_path.endswith('.xml'):
                content = self._format_xml(content)
            elif file_path.endswith('.html'):
                content = self._format_html(content)
            
            return {
                'content': content,
                'metadata': self.get_metadata(file_path),
                'chunks': self._chunk_text(content),
            }
        except Exception as e:
            logger.error(f"Error parsing text file {file_path}: {e}")
            return {'error': str(e)}
    
    def _format_json(self, content: str) -> str:
        """Format JSON content for better readability"""
        try:
            data = json.loads(content)
            return json.dumps(data, indent=2)
        except:
            return content
    
    def _format_xml(self, content: str) -> str:
        """Format XML content using BeautifulSoup"""
        try:
            soup = BeautifulSoup(content, 'xml')
            return soup.prettify()
        except:
            return content
    
    def _format_html(self, content: str) -> str:
        """Extract text from HTML"""
        try:
            soup = BeautifulSoup(content, 'html.parser')
            return soup.get_text(separator='\n', strip=True)
        except:
            return content
    
    def _chunk_text(self, content: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into chunks"""
        if len(content) <= chunk_size:
            return [content]
        
        chunks = []
        start = 0
        while start < len(content):
            end = start + chunk_size
            if end >= len(content):
                chunks.append(content[start:])
                break
            
            # Try to break at sentence boundary
            chunk = content[start:end]
            last_sentence = chunk.rfind('.')
            if last_sentence > chunk_size // 2:
                chunk = chunk[:last_sentence + 1]
                end = start + last_sentence + 1
            
            chunks.append(chunk)
            start = end - overlap
        
        return chunks


class PDFParser(BaseParser):
    """Parser for PDF files"""

    def __init__(self, config: Dict[str, Any] = None):
        super().__init__(config)
        api_key = os.environ.get("MISTRAL_API_KEY")
        if not api_key:
            raise RuntimeError("MISTRAL_API_KEY environment variable not set")
        self.client = Mistral(api_key=api_key)
    
    def can_parse(self, file_path: str, mime_type: str) -> bool:
        pdf_file = Path(file_path)
        return (file_path.lower().endswith('.pdf') or mime_type == 'application/pdf') and pdf_file.is_file()
    
    def parse(self, file_path: str, **kwargs) -> Dict[str, Any]:
        try:
            pdf_file = Path(file_path)
            uploaded_file = self.client.files.upload(
                file={
                    "file_name": pdf_file.stem,
                    "content": pdf_file.read_bytes(),
                },
                purpose="ocr",
            )
            uploaded_file = self.client.files.upload(
                file={
                    "file_name": pdf_file.stem,
                    "content": pdf_file.read_bytes(),
                },
                purpose="ocr",
            )

            signed_url = self.client.files.get_signed_url(file_id=uploaded_file.id, expiry=1)

            ocr_response = self.client.ocr.process(document=DocumentURLChunk(document_url=signed_url.url), model="mistral-ocr-latest", include_image_base64=False)

            content = self.get_combined_markdown(ocr_response)

            logger.info(f"PDF Content {content} parsed successfully with Mistral OCR")

            return {
                'content': content,
                'metadata': self.get_metadata(file_path),
                'chunks': self._chunk_text(content),
            }
        except Exception as e:
            logger.error(f"Error parsing PDF {file_path}: {e}")
            return {'error': str(e)}
        
    def get_combined_markdown(self, ocr_response: OCRResponse) -> str:
        markdowns: list[str] = []
        for page in ocr_response.pages:
            markdowns.append(page.markdown)

        return "\n\n".join(markdowns)
    
    def _extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF using PyPDF2"""
        content = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    content += page.extract_text() + "\n"
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
        return content
    
    # TODO: Implement OCR for PDF files
    def _ocr_pdf(self, file_path: str) -> str:
        """OCR PDF using images"""
        try:
            images = convert_from_path(file_path)
            ocr_texts = []
            for i, image in enumerate(images):
                # Convert to grayscale for better OCR
                if image.mode != 'L':
                    image = image.convert('L')
                image = ImageEnhance.Contrast(image).enhance(1.5)
                image = image.filter(ImageFilter.MedianFilter(size=3))
                text = pytesseract.image_to_string(image, lang='hin+eng')
                if text.strip():
                    ocr_texts.append(f"Page {i+1}:\n{text.strip()}")
            return '\n\n'.join(ocr_texts)
        except Exception as e:
            logger.error(f"OCR error for PDF {file_path}: {e}")
            return ""
    
    def _chunk_text(self, content: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into chunks"""
        text_parser = TextParser()
        return text_parser._chunk_text(content, chunk_size, overlap)


class ImageParser(BaseParser):
    """Parser for image files with OCR capabilities"""
    
    SUPPORTED_EXTENSIONS = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']
    SUPPORTED_MIME_TYPES = ['image/jpeg', 'image/png', 'image/gif', 'image/bmp', 'image/tiff', 'image/webp']
    
    def __init__(self, config: Dict[str, Any] = None):
        super().__init__(config)
        self.ocr_engine = config.get('ocr_engine', 'tesseract') if config else 'tesseract'
        self.easyocr_reader = None
        
        if self.ocr_engine == 'easyocr' and EASYOCR_AVAILABLE:
            self.easyocr_reader = easyocr.Reader(['en'])
    
    def can_parse(self, file_path: str, mime_type: str) -> bool:
        ext = os.path.splitext(file_path)[1].lower()
        return ext in self.SUPPORTED_EXTENSIONS or mime_type in self.SUPPORTED_MIME_TYPES
    
    def parse(self, file_path: str, **kwargs) -> Dict[str, Any]:
        try:
            # Extract text using OCR
            ocr_text = self._extract_text_from_image(file_path)
            
            # Get image metadata
            image_metadata = self._get_image_metadata(file_path)
            
            # Generate image description (placeholder for AI description)
            description = self._generate_image_description(file_path)
            
            combined_content = f"OCR Text:\n{ocr_text}\n\nImage Description:\n{description}"
            
            return {
                'content': combined_content,
                'ocr_text': ocr_text,
                'description': description,
                'metadata': {**self.get_metadata(file_path), **image_metadata},
                'chunks': [combined_content] if combined_content.strip() else [],
            }
        except Exception as e:
            logger.error(f"Error parsing image {file_path}: {e}")
            return {'error': str(e)}
    
    def _extract_text_from_image(self, file_path: str) -> str:
        """Extract text from image using OCR"""
        try:
            if self.ocr_engine == 'easyocr' and self.easyocr_reader:
                result = self.easyocr_reader.readtext(file_path)
                return '\n'.join([text[1] for text in result])
            else:
                # Use Tesseract
                image = Image.open(file_path)

                if image.mode != 'L':
                    image = image.convert('L') # Convert to grayscale

                # Enhance image for better OCR
                image = ImageEnhance.Contrast(image).enhance(1.5)
                image = image.filter(ImageFilter.MedianFilter(size=3))

                # OCR with optimized settings
                return pytesseract.image_to_string(
                        image, 
                        lang='hin+eng', 
                        config='--psm 6 --oem 3 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz .,!?;:'
                    ).strip()
        except Exception as e:
            logger.error(f"Error extracting text from image: {e}")
            return ""
    
    def _get_image_metadata(self, file_path: str) -> Dict[str, Any]:
        """Get image metadata"""
        try:
            with Image.open(file_path) as img:
                return {
                    'width': img.width,
                    'height': img.height,
                    'format': img.format,
                    'mode': img.mode,
                }
        except Exception as e:
            logger.error(f"Error getting image metadata: {e}")
            return {}
    
    def _generate_image_description(self, file_path: str) -> str:
        """Generate image description (placeholder for AI vision model)"""
        # This would integrate with OpenAI Vision API or similar
        return "Image description not available (requires AI vision model integration)"


class AudioParser(BaseParser):
    """Parser for audio files with speech-to-text"""
    
    SUPPORTED_EXTENSIONS = ['.mp3', '.wav', '.m4a', '.ogg', '.flac', '.aac']
    SUPPORTED_MIME_TYPES = ['audio/mpeg', 'audio/wav', 'audio/mp4', 'audio/ogg', 'audio/flac', 'audio/aac']
    
    def __init__(self, config: Dict[str, Any] = None):
        super().__init__(config)
        self.whisper_model = None
        if MEDIA_PROCESSING_AVAILABLE:
            model_name = config.get('whisper_model', 'base') if config else 'base'
            try:
                self.whisper_model = whisper.load_model(model_name)
            except Exception as e:
                logger.error(f"Error loading Whisper model: {e}")
    
    def can_parse(self, file_path: str, mime_type: str) -> bool:
        if not MEDIA_PROCESSING_AVAILABLE:
            return False
        ext = os.path.splitext(file_path)[1].lower()
        return ext in self.SUPPORTED_EXTENSIONS or mime_type in self.SUPPORTED_MIME_TYPES
    
    def parse(self, file_path: str, **kwargs) -> Dict[str, Any]:
        try:
            # Convert audio to WAV if needed
            wav_path = self._convert_to_wav(file_path)
            
            # Transcribe audio
            transcript = self._transcribe_audio(wav_path)
            
            # Clean up temporary file
            if wav_path != file_path:
                os.unlink(wav_path)
            
            return {
                'content': transcript,
                'metadata': self.get_metadata(file_path),
                'chunks': self._chunk_text(transcript),
            }
        except Exception as e:
            logger.error(f"Error parsing audio {file_path}: {e}")
            return {'error': str(e)}
    
    def _convert_to_wav(self, file_path: str) -> str:
        """Convert audio file to WAV format"""
        try:
            audio = AudioSegment.from_file(file_path)
            wav_path = file_path + '_temp.wav'
            audio.export(wav_path, format='wav')
            return wav_path
        except Exception as e:
            logger.error(f"Error converting audio to WAV: {e}")
            return file_path
    
    def _transcribe_audio(self, file_path: str) -> str:
        """Transcribe audio using Whisper"""
        try:
            if self.whisper_model:
                result = self.whisper_model.transcribe(file_path)
                return result['text']
            else:
                return "Audio transcription not available (Whisper model not loaded)"
        except Exception as e:
            logger.error(f"Error transcribing audio: {e}")
            return f"Error transcribing audio: {e}"
    
    def _chunk_text(self, content: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into chunks"""
        text_parser = TextParser()
        return text_parser._chunk_text(content, chunk_size, overlap)


class VideoParser(BaseParser):
    """Parser for video files with audio extraction and frame analysis"""
    
    SUPPORTED_EXTENSIONS = ['.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm']
    SUPPORTED_MIME_TYPES = ['video/mp4', 'video/avi', 'video/quicktime', 'video/x-msvideo', 'video/webm']
    
    def __init__(self, config: Dict[str, Any] = None):
        super().__init__(config)
        self.whisper_model = None
        self.extract_frames = config.get('extract_frames', True) if config else True
        self.frame_interval = config.get('frame_interval', 30) if config else 30  # seconds
        
        if MEDIA_PROCESSING_AVAILABLE:
            model_name = config.get('whisper_model', 'base') if config else 'base'
            try:
                self.whisper_model = whisper.load_model(model_name)
            except Exception as e:
                logger.error(f"Error loading Whisper model: {e}")
    
    def can_parse(self, file_path: str, mime_type: str) -> bool:
        if not MEDIA_PROCESSING_AVAILABLE:
            return False
        ext = os.path.splitext(file_path)[1].lower()
        return ext in self.SUPPORTED_EXTENSIONS or mime_type in self.SUPPORTED_MIME_TYPES
    
    def parse(self, file_path: str, **kwargs) -> Dict[str, Any]:
        try:
            # Extract audio from video
            audio_path = self._extract_audio(file_path)
            
            # Transcribe audio
            transcript = self._transcribe_audio(audio_path)
            
            # Extract frames for OCR
            frame_text = ""
            if self.extract_frames:
                frame_text = self._extract_text_from_frames(file_path)
            
            # Combine content
            combined_content = f"Audio Transcript:\n{transcript}\n\nFrame Text:\n{frame_text}"
            
            # Clean up temporary files
            if os.path.exists(audio_path):
                os.unlink(audio_path)
            
            return {
                'content': combined_content,
                'transcript': transcript,
                'frame_text': frame_text,
                'metadata': self.get_metadata(file_path),
                'chunks': self._chunk_text(combined_content),
            }
        except Exception as e:
            logger.error(f"Error parsing video {file_path}: {e}")
            return {'error': str(e)}
    
    def _extract_audio(self, file_path: str) -> str:
        """Extract audio from video"""
        try:
            video = AudioSegment.from_file(file_path)
            audio_path = file_path + '_temp_audio.wav'
            video.export(audio_path, format='wav')
            return audio_path
        except Exception as e:
            logger.error(f"Error extracting audio from video: {e}")
            return ""
    
    def _transcribe_audio(self, audio_path: str) -> str:
        """Transcribe audio using Whisper"""
        try:
            if self.whisper_model and os.path.exists(audio_path):
                result = self.whisper_model.transcribe(audio_path)
                return result['text']
            else:
                return "Audio transcription not available"
        except Exception as e:
            logger.error(f"Error transcribing audio: {e}")
            return f"Error transcribing audio: {e}"
    
    def _extract_text_from_frames(self, file_path: str) -> str:
        """Extract text from video frames using OCR"""
        try:
            cap = cv2.VideoCapture(file_path)
            fps = cap.get(cv2.CAP_PROP_FPS)
            frame_interval_frames = int(fps * self.frame_interval)
            
            frame_texts = []
            frame_count = 0
            
            while True:
                ret, frame = cap.read()
                if not ret:
                    break
                
                if frame_count % frame_interval_frames == 0:
                    # Convert frame to PIL Image
                    frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                    pil_image = Image.fromarray(frame_rgb)
                    
                    # Extract text using OCR
                    text = pytesseract.image_to_string(pil_image)
                    if text.strip():
                        frame_texts.append(f"Frame {frame_count//frame_interval_frames}: {text.strip()}")
                
                frame_count += 1
            
            cap.release()
            return '\n'.join(frame_texts)
        except Exception as e:
            logger.error(f"Error extracting text from frames: {e}")
            return ""
    
    def _chunk_text(self, content: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into chunks"""
        text_parser = TextParser()
        return text_parser._chunk_text(content, chunk_size, overlap)


class DocumentParser(BaseParser):
    """Parser for Office documents (Word, Excel, PowerPoint)"""
    
    SUPPORTED_EXTENSIONS = ['.docx', '.xlsx', '.pptx']
    SUPPORTED_MIME_TYPES = [
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    ]
    
    def can_parse(self, file_path: str, mime_type: str) -> bool:
        ext = os.path.splitext(file_path)[1].lower()
        return ext in self.SUPPORTED_EXTENSIONS or mime_type in self.SUPPORTED_MIME_TYPES
    
    def parse(self, file_path: str, **kwargs) -> Dict[str, Any]:
        try:
            ext = os.path.splitext(file_path)[1].lower()
            logger.info(f"DocumentParser: Starting parse for {file_path} with extension {ext}")
            
            if ext == '.docx':
                logger.debug(f"DocumentParser: Parsing DOCX file {file_path}")
                result = self._parse_docx(file_path)
            elif ext == '.xlsx':
                logger.debug(f"DocumentParser: Parsing XLSX file {file_path}")
                result = self._parse_xlsx(file_path)
            elif ext == '.pptx':
                logger.debug(f"DocumentParser: Parsing PPTX file {file_path}")
                result = self._parse_pptx(file_path)
            else:
                logger.warning(f"DocumentParser: Unsupported file type {ext} for file {file_path}")
                result = {'error': f'Unsupported file type: {ext}'}
            
            logger.info(f"DocumentParser: Finished parsing {file_path}")
            return result
        except Exception as e:
            logger.error(f"Error parsing document {file_path}: {e}")
            return {'error': str(e)}
    
    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """Parse Word document"""
        doc = DocxDocument(file_path)
        content = '\n'.join([para.text for para in doc.paragraphs])
        
        return {
            'content': content,
            'metadata': self.get_metadata(file_path),
            'chunks': self._chunk_text(content),
        }
    
    def _parse_xlsx(self, file_path: str) -> Dict[str, Any]:
        """Parse Excel spreadsheet"""
        workbook = openpyxl.load_workbook(file_path)
        content_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_content = f"Sheet: {sheet_name}\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                if row_text.strip():
                    sheet_content += row_text + '\n'
            
            content_parts.append(sheet_content)
        
        content = '\n'.join(content_parts)
        
        return {
            'content': content,
            'metadata': self.get_metadata(file_path),
            'chunks': self._chunk_text(content),
        }
    
    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """Parse PowerPoint presentation"""
        prs = Presentation(file_path)
        content_parts = []
        
        for i, slide in enumerate(prs.slides):
            slide_content = f"Slide {i+1}:\n"
            
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    slide_content += shape.text + '\n'
            
            content_parts.append(slide_content)
        
        content = '\n'.join(content_parts)
        
        return {
            'content': content,
            'metadata': self.get_metadata(file_path),
            'chunks': self._chunk_text(content),
        }
    
    def _chunk_text(self, content: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into chunks"""
        text_parser = TextParser()
        return text_parser._chunk_text(content, chunk_size, overlap)


class ParserRegistry:
    """Registry for managing document parsers"""
    
    def __init__(self):
        self.parsers = []
        self._register_default_parsers()
    
    def _register_default_parsers(self):
        """Register default parsers"""
        self.register_parser(TextParser())
        self.register_parser(PDFParser())
        self.register_parser(ImageParser())
        self.register_parser(DocumentParser())
        
        if MEDIA_PROCESSING_AVAILABLE:
            self.register_parser(AudioParser())
            self.register_parser(VideoParser())
    
    def register_parser(self, parser: BaseParser):
        """Register a new parser"""
        self.parsers.append(parser)
    
    def get_parser(self, file_path: str, mime_type: str) -> Optional[BaseParser]:
        """Get the appropriate parser for a file"""
        for parser in self.parsers:
            if parser.can_parse(file_path, mime_type):
                return parser
        return None
    
    def get_parsers_for_type(self, file_type: str) -> List[BaseParser]:
        """Get all parsers that can handle a specific file type"""
        parsers = []
        for parser in self.parsers:
            if parser.can_parse(f"test.{file_type}", ""):
                parsers.append(parser)
        return parsers


# Global parser registry
parser_registry = ParserRegistry()

def encode_image(image_path):
    """Encode the image to base64."""
    try:
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')
    except FileNotFoundError:
        print(f"Error: The file {image_path} was not found.")
        return None
    except Exception as e:  # Added general exception handling
        print(f"Error: {e}")
        return None